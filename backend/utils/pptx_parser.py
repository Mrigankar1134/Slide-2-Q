from typing import List
from pptx import Presentation


def extract_slides_text(file_path: str) -> List[str]:
    prs = Presentation(file_path)
    slide_texts: List[str] = []
    for slide in prs.slides:
        texts: List[str] = []
        # Title
        if slide.shapes.title and slide.shapes.title.text:
            texts.append(slide.shapes.title.text)
        # All text-containing shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        joined = "\n".join(texts).strip()
        slide_texts.append(joined)
    return slide_texts
